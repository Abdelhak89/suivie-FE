import sys
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

# ✅ placeholders dont la valeur doit être en taille 8
FONT8_KEYS = {
    "{{DESCRIPTION_FE}}",
    "{{PARTICIPANTS}}",
    "{{QUOI_COMBIEN}}",
    "{{QUI}}",
    "{{OU}}",
    "{{QUAND}}",
    "{{RECURRENCE}}",
}

def replace_in_textframe(tf, mapping):
    tf.word_wrap = True
    # ✅ shrink-to-fit la zone si trop long
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for p in tf.paragraphs:
        for run in p.runs:
            txt = run.text
            for k, v in mapping.items():
                if k in txt:
                    run.text = txt.replace(k, v)

                    # ✅ force font-size 8 sur les zones data
                    if k in FONT8_KEYS:
                        run.font.size = Pt(8)

def replace_in_shape(shape, mapping):
    # TextFrame
    if shape.has_text_frame:
        replace_in_textframe(shape.text_frame, mapping)

    # Tables
    if shape.has_table:
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                replace_in_textframe(cell.text_frame, mapping)

def main():
    if len(sys.argv) < 4:
        print("Usage: fill_clinique_qualite.py <template.pptx> <payload.json> <out.pptx>", file=sys.stderr)
        sys.exit(1)

    template_path = sys.argv[1]
    payload_path = sys.argv[2]
    out_path = sys.argv[3]

    with open(payload_path, "r", encoding="utf-8") as f:
        payload = json.load(f)

    mapping = payload.get("mapping", {})
    mapping = {k: ("" if v is None else str(v)) for k, v in mapping.items()}

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, mapping)

    prs.save(out_path)

if __name__ == "__main__":
    main()
